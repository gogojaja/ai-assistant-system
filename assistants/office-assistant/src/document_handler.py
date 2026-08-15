"""
模块名称：document_handler
功能描述：文档处理（Word/Excel/PPT 文件下载、分析、摘要生成、PPT 生成、回复组装）
        文档处理能力已弱化：PPT 生成优先委托外部工具（PPTAgent/dashi-ppt），摘要优先委托后端 LLM，本地 core 仅作兜底
对外接口：
    - process_document_file(file_key, message_id, open_id, file_name): 处理文档消息
    - process_office_text(cmd_text, open_id, target_id, receive_id_type): 处理 #2/#office 文本命令
依赖：
    - 标准库：logging, tempfile, os, sys, pathlib, re
    - 第三方：无（requests/yaml 由 shared.backend_utils 托管）
    - 项目内：shared.feishu_api (download_file, send_message, send_file_message),
               shared.backend_utils (call_api, get_backend_config, clean_reply),
               core.word_processor (WordProcessor, 兜底),
               core.summarizer (DocumentSummarizer, 兜底),
               core.excel_processor (ExcelProcessor, 兜底),
               core.ppt_generator (generate_presentation, generate_from_text, 兜底),
               external_doc_tools (generate_ppt_via_pptagent, generate_ppt_via_dashi, summarize_doc_via_external)
版本：v3.1
更新记录：
    - 2026-08-16: 文档处理能力弱化，PPT 生成优先调 PPTAgent/dashi-ppt，摘要优先委托外部 LLM，本地 core 降为兜底
    - 2026-05-26: 新增 PPT 文件支持、#2/#office 文本命令、PPT 生成功能
    - 2026-05-23: 初始创建，从 callback_server.py 剥离文档处理逻辑，包含 Excel 智能摘要
    - 2026-05-25: Excel 摘要改用独立 API 调用，不再依赖 main.talk
    - 2026-05-26: 抽取重复代码到 shared/backend_utils.py；启动文件夹监控
"""
import logging
import re
import tempfile
import os
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).parent.parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))
sys.path.insert(0, str(PROJECT_ROOT / "assistants/chat-assistant/src"))
sys.path.insert(0, str(PROJECT_ROOT / "assistants/office-assistant/src"))

from shared.feishu_api import download_file, send_message, send_file_message
from shared.backend_utils import call_api

logger = logging.getLogger(__name__)

try:
    from core.word_processor import WordProcessor
    from core.summarizer import DocumentSummarizer
    from core.excel_processor import ExcelProcessor
    from core.ppt_generator import generate_presentation, generate_from_text
    DOC_PROCESSING_AVAILABLE = True
    doc_summarizer = DocumentSummarizer()
except ImportError as e:
    DOC_PROCESSING_AVAILABLE = False
    doc_summarizer = None
    logger.error(f"文档处理模块导入失败: {e}")

# 外部文档工具适配层（PPTAgent / dashi-ppt / 外部摘要）
try:
    from external_doc_tools import (
        check_external_ready,
        generate_ppt_via_pptagent,
        generate_ppt_via_dashi,
        summarize_doc_via_external,
    )
    EXTERNAL_TOOLS_AVAILABLE = True
except ImportError as e:
    EXTERNAL_TOOLS_AVAILABLE = False
    logger.error(f"外部文档工具适配层导入失败: {e}")

# 启动办公文件夹监控
try:
    from core.folder_monitor import start_monitor
    _watch_dir = str(PROJECT_ROOT / "data" / "office")
    start_monitor(_watch_dir)
    logger.info(f"办公文件夹监控已启动: {_watch_dir}")
except Exception as e:
    logger.warning(f"办公文件夹监控启动失败: {e}")


def _is_valid_summary(text: str) -> bool:
    if not text or len(text) < 10:
        return False
    if text == "（摘要为空）":
        return False
    if text.startswith("（"):
        return False
    return True


def _build_fallback_summary(structure_hint: str, data_text: str) -> str:
    """当 AI 摘要失败时，生成基于规则的摘要（不含 structure_hint，避免重复）"""
    lines = data_text.strip().split("\n")
    header_line = ""
    data_rows = []
    for line in lines:
        if "\t" in line:
            cols = [c.strip() for c in line.split("\t")]
            if not header_line:
                header_line = cols
            else:
                data_rows.append(cols)
    parts = []
    if header_line:
        valid = [h for h in header_line if h]
        if valid:
            parts.append(f"列字段：{'、'.join(valid[:10])}")
    if data_rows:
        parts.append(f"共 {len(data_rows)} 条记录")
        sample = []
        for row in data_rows[:3]:
            vals = [v for v in row[:5] if v]
            if vals:
                sample.append("、".join(vals[:3]))
        if sample:
            parts.append("例如：" + "；".join(sample))
    return "；".join(parts)


def generate_excel_summary(data_text: str, structure_hint: str = "") -> tuple:
    """
    生成 Excel 数据摘要，返回 (summary_text, is_ai)。
    外部委托优先，本地 call_api 兜底，最后规则兜底。
    is_ai=True 表示 AI 生成，False 表示规则兜底。
    """
    stripped = data_text.strip()
    if not stripped or len(stripped) < 10:
        return ("", False)
    fallback = _build_fallback_summary(structure_hint, stripped)
    try:
        lines = stripped.split("\n")
        header_line = ""
        data_rows = []
        for line in lines:
            if "\t" in line:
                cols = [c.strip() for c in line.split("\t")]
                if not header_line:
                    header_line = cols
                else:
                    data_rows.append(cols)
        samples = []
        for row in data_rows[:5]:
            vals = [v for v in row[:6] if v]
            if vals:
                samples.append("、".join(vals[:4]))
        concise_parts = [f"列字段：{'、'.join(header_line[:8])}" if header_line else ""]
        concise_parts.append(f"共 {len(data_rows)} 条数据记录")
        if samples:
            concise_parts.append(f"样例数据：{'；'.join(samples)}")
        concise = "；".join(p for p in concise_parts if p)
        # 外部摘要优先
        if EXTERNAL_TOOLS_AVAILABLE:
            try:
                ext_summary, ext_ok = summarize_doc_via_external(
                    stripped,
                    hint=f"Excel 数据。数据概况：{concise}",
                )
                if ext_ok:
                    return (ext_summary, True)
            except Exception:
                pass
        # 本地 call_api 兜底
        hint = f"文件结构：{structure_hint}\n\n" if structure_hint else ""
        user_msg = (
            f"{hint}"
            f"数据概况：{concise}\n\n"
            f"这是一份Excel数据，请用1-3句话简要总结核心信息："
        )
        result = call_api([{"role": "user", "content": user_msg}])
        if result and _is_valid_summary(result):
            return (result, True)
    except Exception:
        pass
    return (fallback, False)


_SIZE_LIMITS = {
    '.docx': 30 * 1024 * 1024,
    '.xlsx': 15 * 1024 * 1024,
    '.pptx': 30 * 1024 * 1024,
}


def _extract_pptx_text(file_path: str) -> dict:
    """提取 PPTX 文件的文本内容，返回 {text, slide_count, slides}"""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    total_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_content.append(t)
        slide_text = "\n".join(slide_content)
        slides_text.append({"index": i, "text": slide_text})
        total_text.append(slide_text)
    return {
        "text": "\n\n".join(total_text),
        "slide_count": len(prs.slides),
        "slides": slides_text,
    }


# 文档内容缓存：open_id → {"text": str, "title": str}
# 用于 转PPT 命令将上次分析的文档转 PPT
_doc_cache: dict = {}


def _cache_doc_text(open_id: str, text: str, title: str = "文档内容"):
    """缓存用户最近一次文档的文本内容"""
    _doc_cache[open_id] = {"text": text, "title": title}
    if len(_doc_cache) > 100:
        _doc_cache.clear()


def _generate_ppt(ppt_text: str, output_path: str, attachments: list = None) -> tuple:
    """
    统一 PPT 生成入口：优先外部工具（PPTAgent → dashi-ppt），本地 core 兜底。
    返回 (success, result)：success=True 时 result 为输出路径，否则为错误信息。
    """
    if not ppt_text or not ppt_text.strip():
        return (False, "文案为空")

    # 1) 优先 PPTAgent（专业级）
    if EXTERNAL_TOOLS_AVAILABLE:
        try:
            ok, res = generate_ppt_via_pptagent(
                ppt_text, output_path, attachments=attachments, lang="zh"
            )
            if ok:
                return (True, res)
            logger.warning(f"PPTAgent 生成未成功，切换兜底: {res}")
        except Exception as e:
            logger.warning(f"PPTAgent 调用异常，切换兜底: {e}")

        # 2) 其次 dashi-ppt
        try:
            ok, res = generate_ppt_via_dashi(ppt_text, output_path)
            if ok:
                return (True, res)
            logger.warning(f"dashi-ppt 生成未成功，切换本地兜底: {res}")
        except Exception as e:
            logger.warning(f"dashi-ppt 调用异常，切换本地兜底: {e}")

    # 3) 本地 core 兜底
    try:
        if not DOC_PROCESSING_AVAILABLE:
            return (False, "文档处理服务未就绪")
        from core.ppt_generator import generate_via_pandoc, generate_from_text
        result = generate_via_pandoc(ppt_text, output_path)
        if not result:
            result = generate_from_text(ppt_text, output_path)
        if result and os.path.getsize(result) > 1000:
            return (True, result)
        return (False, "本地 PPT 生成失败")
    except Exception as e:
        logger.error(f"本地 PPT 生成异常: {e}", exc_info=True)
        return (False, f"本地 PPT 生成异常：{e}")


_OFFICE_HELP = """
📋 2号AI 办公助理 — 可用命令：
  #办公 help            — 显示此帮助
  #办公 ppt <文案>       — 根据文案生成 PPT
  #办公 生成ppt：<文案>   — 同上
  #办公 转PPT            — 将上次分析的文档转为 PPT（或直接发 转PPT）

📎 发送文件即可分析：
  .docx → 提取摘要 + 支持 转PPT
  .xlsx → 数据分析 + AI 摘要
  .pptx → 提取幻灯片内容 + 支持 转PPT
""".strip()


def process_office_text(cmd_text: str, open_id: str, target_id: str = "", receive_id_type: str = "open_id"):
    """处理 #2/#office 文本命令"""
    cmd_text = cmd_text.strip()
    if not cmd_text or cmd_text == "help":
        send_message(target_id or open_id, _OFFICE_HELP, receive_id_type=receive_id_type)
        return

    # PPT 生成：ppt <文案> 或 生成ppt：<文案>
    ppt_text = ""
    if cmd_text.startswith("ppt ") or cmd_text.startswith("ppt："):
        ppt_text = cmd_text[4:].strip()
    elif cmd_text.startswith("ppt:") or cmd_text.startswith("ppt:"):
        ppt_text = cmd_text[4:].strip()
    elif cmd_text.startswith("生成ppt") or cmd_text.startswith("生成PPT"):
        ppt_text = re.sub(r'^(生成ppt|生成PPT)[：:]\s*', '', cmd_text).strip()
    elif cmd_text.startswith("生成演示"):
        ppt_text = re.sub(r'^生成演示[：:]\s*', '', cmd_text).strip()

    if ppt_text:
        send_message(target_id or open_id, "📊 正在生成 PPT（外部工具优先，预计 1-10 分钟）…", receive_id_type=receive_id_type)
        try:
            output_dir = PROJECT_ROOT / "data" / "office"
            output_dir.mkdir(parents=True, exist_ok=True)
            ts = int(__import__('time').time())
            output_path = str(output_dir / f"presentation_{ts}.pptx")
            ok, result = _generate_ppt(ppt_text, output_path)
            if ok and os.path.getsize(result) > 1000:
                send_ok = send_file_message(target_id or open_id, result, "专业演示文稿.pptx", receive_id_type=receive_id_type)
                if not send_ok:
                    send_message(target_id or open_id, "PPT 已生成，但文件发送失败，请稍后重试。", receive_id_type=receive_id_type)
                else:
                    send_message(target_id or open_id, "✅ PPT 已发送！", receive_id_type=receive_id_type)
            else:
                send_message(target_id or open_id, f"PPT 生成失败：{result if isinstance(result, str) else '请提供更详细的文案'}。可使用空行分隔幻灯片，'-' 标记要点。", receive_id_type=receive_id_type)
        except Exception as e:
            logger.error(f"PPT 生成异常: {e}", exc_info=True)
            send_message(target_id or open_id, f"PPT 生成失败：{str(e)}", receive_id_type=receive_id_type)
        return

    # toppt / 转ppt：从缓存的上次文档生成 PPT
    if cmd_text in ("转PPT",):
        cached = _doc_cache.get(open_id)
        if not cached or not cached["text"].strip():
            send_message(target_id or open_id, "没有找到已分析的文档。请先发送 .docx 或 .pptx 文件。", receive_id_type=receive_id_type)
            return
        ppt_text = cached["text"]
        ppt_title = cached["title"]
        send_message(target_id or open_id, f"📊 正在将「{ppt_title}」转为 PPT（外部工具优先，预计 1-10 分钟）…", receive_id_type=receive_id_type)
        try:
            output_dir = PROJECT_ROOT / "data" / "office"
            output_dir.mkdir(parents=True, exist_ok=True)
            ts = int(__import__('time').time())
            output_path = str(output_dir / f"presentation_{ts}.pptx")
            ok, result = _generate_ppt(ppt_text, output_path)
            if ok and os.path.getsize(result) > 1000:
                send_ok = send_file_message(target_id or open_id, result, f"{ppt_title}.pptx", receive_id_type=receive_id_type)
                if not send_ok:
                    send_message(target_id or open_id, "PPT 已生成，但文件发送失败，请稍后重试。", receive_id_type=receive_id_type)
                else:
                    send_message(target_id or open_id, "✅ PPT 已发送！", receive_id_type=receive_id_type)
            else:
                send_message(target_id or open_id, f"PPT 生成失败：{result if isinstance(result, str) else '文档内容可能不完整'}。", receive_id_type=receive_id_type)
        except Exception as e:
            logger.error(f"toppt 异常: {e}", exc_info=True)
            send_message(target_id or open_id, f"PPT 生成失败：{str(e)}", receive_id_type=receive_id_type)
        return

    send_message(target_id or open_id, f"未知命令：{cmd_text}\n\n{_OFFICE_HELP}", receive_id_type=receive_id_type)


def process_document_file(file_key: str, message_id: str, open_id: str, file_name: str = "document.docx"):
    """处理 Word / Excel / PPT 文件消息"""
    suffix = Path(file_name).suffix.lower()
    if suffix not in ['.docx', '.xlsx', '.pptx']:
        send_message(open_id, f"暂不支持的文档格式：{suffix}")
        return
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        file_path = tmp.name
    if not download_file(message_id, file_key, file_path):
        send_message(open_id, "文件下载失败，请稍后重试。")
        return

    file_size = os.path.getsize(file_path)
    file_size_mb = file_size / (1024 * 1024)
    size_limit = _SIZE_LIMITS.get(suffix, 20 * 1024 * 1024)
    size_limit_mb = size_limit / (1024 * 1024)

    if file_size > size_limit:
        send_message(open_id, f"⚠️ 文件过大（{file_size_mb:.1f}MB），当前限制最大 {size_limit_mb:.0f}MB。请压缩后重试。")
        os.unlink(file_path)
        return

    icon_map = {'.docx': "📄", '.xlsx': "📊", '.pptx': "📽️"}
    label_map = {'.docx': "Word 文件", '.xlsx': "Excel 文件", '.pptx': "PPT 文件"}
    icon = icon_map.get(suffix, "📎")
    type_label = label_map.get(suffix, "文件")
    send_message(open_id, f"{icon} 收到 {type_label}：{file_name}（{file_size_mb:.1f}MB），正在处理…")

    try:
        if suffix == '.docx':
            if not DOC_PROCESSING_AVAILABLE:
                send_message(open_id, "文档处理服务未就绪，请稍后重试。")
                return
            wp = WordProcessor(file_path)
            info = wp.get_summary_info()
            text = wp.extract_text()
            titles = wp.extract_titles()
            _cache_doc_text(open_id, text, info['file_name'])
            summary = ""
            summary_error = ""
            if EXTERNAL_TOOLS_AVAILABLE:
                try:
                    ext_summary, ext_ok = summarize_doc_via_external(
                        text, hint=f"《{info['file_name']}》Word 文档"
                    )
                    if ext_ok:
                        summary = ext_summary
                except Exception as e:
                    logger.warning(f"外部摘要失败，本地兜底: {e}")
            if not summary and DOC_PROCESSING_AVAILABLE:
                result = doc_summarizer.summarize(text, max_points=5)
                summary = result['summary']
                summary_error = result.get('error')
            reply = f"📄 **{info['file_name']}** ({file_size_mb:.1f}MB)\n\n{summary}\n\n"
            reply += f"📊 文档信息：{info['paragraph_count']} 段落, {info['table_count']} 表格, {info['character_count']} 字符"
            if titles:
                reply += f"\n📑 标题：{', '.join([t['text'] for t in titles[:5]])}"
            if summary_error:
                reply += f"\n⚠️ 注意：{summary_error}"
            reply += "\n\n💡 如需转为 PPT，发送：转PPT"

        elif suffix == '.xlsx':
            if not DOC_PROCESSING_AVAILABLE:
                send_message(open_id, "文档处理服务未就绪，请稍后重试。")
                return
            processor = ExcelProcessor(file_path)
            structure = processor.analyze()
            reply = f"{structure}\n（{file_size_mb:.1f}MB）"
            data_text = processor.get_data_text()
            if data_text:
                summary, is_ai = generate_excel_summary(data_text, structure_hint=structure)
                if summary:
                    header = "\n\n🤖 AI 摘要：\n" if is_ai else "\n\n📋 数据概况：\n"
                    max_len = 4900
                    base_len = len(reply) + len(header)
                    if base_len + len(summary) > max_len:
                        available = max_len - base_len
                        if available > 10:
                            truncated = summary[:available]
                            last_sentence = max(truncated.rfind("。"), truncated.rfind("，"))
                            if last_sentence > available // 2:
                                summary = truncated[:last_sentence+1] + "…"
                            else:
                                summary = truncated[:available-1] + "…"
                        else:
                            summary = "（摘要过长，已省略）"
                    reply = reply + header + summary

        else:
            ppt_data = _extract_pptx_text(file_path)
            slide_count = ppt_data["slide_count"]
            full_text = ppt_data["text"]
            _cache_doc_text(open_id, full_text, file_name)
            reply = f"📽️ **{file_name}** ({file_size_mb:.1f}MB)\n\n共 {slide_count} 页幻灯片"
            if full_text:
                truncated = full_text[:2000]
                reply += f"\n\n📝 内容预览：\n{truncated}"
                if len(full_text) > 2000:
                    reply += "\n…（内容较长，已截取前 2000 字符）"
            reply += "\n\n💡 如需转为 PPT，发送：转PPT"

        send_message(open_id, reply)
    except Exception as e:
        logger.error(f"文档处理异常: {e}", exc_info=True)
        send_message(open_id, f"文档处理失败：{str(e)}")
    finally:
        try:
            os.unlink(file_path)
        except Exception:
            pass
