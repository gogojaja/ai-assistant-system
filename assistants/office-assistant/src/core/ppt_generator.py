"""
模块名称：ppt_generator
功能描述：专业级 PPT 生成器，支持 python-pptx 引擎和 Pandoc 管道
对外接口：
    - generate_presentation(title, slides, output_path): 根据结构化数据生成 PPTX
    - generate_from_text(text, output_path): 自动解析文本生成 PPTX
    - check_pandoc(): 检测 Pandoc 是否可用
依赖：
    - 标准库：logging, pathlib, re, subprocess, tempfile
    - 第三方：pptx
    - 项目内：无
版本：v2.0
更新记录：
    - 2026-05-26: 全面重写，专业级模板引擎 + Pandoc 管道支持
"""
import logging
import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ============================================================
# 专业级 PPT 设计系统
# ============================================================

# --- 色彩系统 ---
C_PRIMARY = RGBColor(0x1a, 0x27, 0x44)        # 深海军蓝
C_ACCENT = RGBColor(0x2b, 0x5f, 0xc3)         # 活力蓝
C_ACCENT_LIGHT = RGBColor(0x4a, 0x7f, 0xd5)   # 浅蓝
C_BG_LIGHT = RGBColor(0xf5, 0xf7, 0xfa)       # 浅灰背景
C_BG_GRADIENT_END = RGBColor(0xe8, 0xed, 0xf5)  # 渐变终点
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1a, 0x1a, 0x2e)            # 正文色
C_TEXT_SECONDARY = RGBColor(0x5a, 0x6a, 0x7a)  # 辅助文本
C_LINE = RGBColor(0x2b, 0x5f, 0xc3)           # 装饰线
C_HIGHLIGHT = RGBColor(0xc0, 0x39, 0x2b)      # 强调色

# --- 字体系统 ---
FONT_TITLE = "Helvetica"
FONT_BODY = "Helvetica"

# --- 尺寸系统 (16:9 宽屏) ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = SLIDE_W - 2 * MARGIN
TITLE_TOP = Inches(0.5)
TITLE_H = Inches(1.0)
CONTENT_TOP = Inches(1.8)
CONTENT_H = SLIDE_H - CONTENT_TOP - Inches(0.5)

# --- 字号系统 ---
SIZE_TITLE_MAIN = Pt(42)
SIZE_TITLE_SUB = Pt(20)
SIZE_SECTION = Pt(40)
SIZE_CONTENT_TITLE = Pt(28)
SIZE_BODY = Pt(18)
SIZE_BULLET = Pt(16)
SIZE_FOOTER = Pt(10)
SIZE_CAPTION = Pt(14)

SHAPES_STYLE = {
    'accent_bar': {'width': Inches(0.06), 'color': C_LINE},
    'header_bg': {'color': C_PRIMARY},
    'accent_dot': {'size': Inches(0.12), 'color': C_ACCENT},
}


# ============================================================
# 辅助渲染函数
# ============================================================

def _add_shape(slide, left, top, width, height, fill_color=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def _add_textbox(slide, left, top, width, height, text="", font_size=SIZE_BODY,
                 font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    """添加文本框并设置样式"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_text(text_frame, items, font_size=SIZE_BULLET, font_color=C_TEXT,
                     bullet_char="•", space_after=Pt(6)):
    """向文本框添加项目符号列表"""
    tf = text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.name = FONT_BODY
        p.space_after = space_after
        p.level = 0
    return text_frame


def _add_footer(slide, text="", page_num=0):
    """添加幻灯片页脚"""
    fg_color = C_TEXT_SECONDARY
    footer_text = text
    if page_num > 0:
        if footer_text:
            footer_text += f"    |    {page_num}"
        else:
            footer_text = str(page_num)
    _add_textbox(slide, MARGIN, SLIDE_H - Inches(0.4), CONTENT_W, Inches(0.3),
                 footer_text, SIZE_FOOTER, fg_color, alignment=PP_ALIGN.CENTER)


# ============================================================
# 幻灯片创建器
# ============================================================

def _add_title_slide(prs, title_text, subtitle_text=""):
    """创建封面幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色块：上半部分
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, C_PRIMARY)

    # 装饰线
    _add_shape(slide, MARGIN, Inches(3.2), Inches(2.5), Inches(0.04), C_ACCENT)

    # 标题
    _add_textbox(slide, MARGIN, Inches(3.5), CONTENT_W, Inches(1.5),
                 title_text, SIZE_TITLE_MAIN, C_WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # 副标题
    if subtitle_text:
        _add_textbox(slide, MARGIN, Inches(5.0), CONTENT_W, Inches(0.8),
                     subtitle_text, SIZE_TITLE_SUB, C_ACCENT_LIGHT,
                     alignment=PP_ALIGN.LEFT)

    return slide


def _add_section_slide(prs, title_text):
    """创建章节分隔幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 深色背景
    _add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, C_PRIMARY)

    # 装饰角标
    _add_shape(slide, MARGIN, Inches(2.8), Inches(0.8), Inches(0.06), C_ACCENT)

    # 章节标题
    _add_textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(1.5),
                 title_text, SIZE_SECTION, C_WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT)

    return slide


def _add_content_slide(prs, title_text, items, notes=""):
    """创建标准内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部标题栏
    _add_shape(slide, 0, 0, SLIDE_W, Inches(1.3), C_PRIMARY)
    _add_textbox(slide, MARGIN, Inches(0.25), CONTENT_W, Inches(0.9),
                 title_text, SIZE_CONTENT_TITLE, C_WHITE, bold=True)

    # 正文区域
    content_box = _add_textbox(slide, MARGIN, CONTENT_TOP, CONTENT_W,
                                Inches(5.0), "", SIZE_BODY)
    tf = content_box.text_frame
    if items:
        _add_bullet_text(tf, items, SIZE_BULLET, C_TEXT)
    else:
        p = tf.paragraphs[0]
        p.text = "（无内容）"
        p.font.size = SIZE_BODY
        p.font.color.rgb = C_TEXT_SECONDARY
        p.font.name = FONT_BODY

    # 底部装饰线
    _add_shape(slide, MARGIN, SLIDE_H - Inches(0.45), CONTENT_W, Inches(0.015),
               C_BG_LIGHT)

    return slide


def _add_two_column_slide(prs, title_text, left_items, right_items):
    """创建双栏内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部标题栏
    _add_shape(slide, 0, 0, SLIDE_W, Inches(1.3), C_PRIMARY)
    _add_textbox(slide, MARGIN, Inches(0.25), CONTENT_W, Inches(0.9),
                 title_text, SIZE_CONTENT_TITLE, C_WHITE, bold=True)

    # 左栏
    col_w = (CONTENT_W - Inches(0.5)) / 2
    left_box = _add_textbox(slide, MARGIN, CONTENT_TOP, col_w, Inches(5.0))
    _add_bullet_text(left_box.text_frame, left_items, SIZE_BULLET, C_TEXT)

    # 分隔线
    sep_x = MARGIN + col_w + Inches(0.2)
    _add_shape(slide, sep_x, CONTENT_TOP, Inches(0.025), Inches(4.5), C_BG_LIGHT)

    # 右栏
    right_box = _add_textbox(slide, sep_x + Inches(0.3), CONTENT_TOP, col_w, Inches(5.0))
    _add_bullet_text(right_box.text_frame, right_items, SIZE_BULLET, C_TEXT)

    return slide


# ============================================================
# 文本解析引擎
# ============================================================

def _parse_text_to_slides(text):
    """将纯文本解析为结构化幻灯片数据

    解析规则：
    - 首行 = 演示标题
    - `## 标题` = 章节分隔页
    - 非 bullet 行 = 幻灯片标题（自动创建新幻灯片）
    - `- / • / *` 开头 = 当前幻灯片要点
    - `左|右` = 双栏内容
    - 空行 = 分隔幻灯片（下一非 bullet 行成为新幻灯片标题）
    """
    lines = [line.strip() for line in text.strip().split("\n")]
    if not lines:
        return None

    title = lines[0]
    slides = []

    # current: 当前正在构建的幻灯片
    # section_mode: True 表示刚经过 ## 章节页，后续行需要分组处理
    current = None
    section_mode = False

    def finalize_current():
        nonlocal current
        if current:
            slides.append(current)
            current = None

    for line in lines[1:]:
        if not line:
            finalize_current()
            continue

        # 章节分隔
        if line.startswith("##") or (line.startswith("===") and line.endswith("===")):
            finalize_current()
            section_title = line.strip("#= ")
            slides.append({"type": "section", "title": section_title})
            section_mode = True
            continue

        # 双栏 (左|右)
        if "|" in line and not line.startswith("-") and not line.startswith("•") and not line.startswith("*"):
            parts = [p.strip() for p in line.split("|")]
            if len(parts) == 2 and parts[0] and parts[1]:
                if current and current.get("type") == "two_column":
                    current["left"].append(parts[0])
                    current["right"].append(parts[1])
                else:
                    finalize_current()
                    current = {"type": "two_column", "title": "对比",
                               "left": [parts[0]], "right": [parts[1]]}
                section_mode = False
                continue

        # Bullet 行
        is_bullet = line.startswith("-") or line.startswith("•") or line.startswith("*")
        clean_text = line.lstrip("-•* ")

        if is_bullet:
            if current is None:
                current = {"type": "content", "title": "要点", "content": [clean_text]}
            elif current["type"] == "content":
                current["content"].append(clean_text)
            section_mode = False

        elif section_mode:
            # 章节模式下：首行为标题，后续行自动成为内容
            if current is None:
                current = {"type": "content", "title": line, "content": []}
            else:
                current["content"].append(line)

        else:
            # 非章节模式下的非 bullet 行：结束当前幻灯片，创建新幻灯片
            finalize_current()
            current = {"type": "content", "title": line, "content": []}

    finalize_current()

    if not slides:
        return None

    return {"title": title, "slides": slides}


# ============================================================
# 核心生成函数
# ============================================================

def generate_presentation(title, slides_data, output_path):
    """根据结构化数据生成专业级 PPTX

    Args:
        title: 演示标题
        slides_data: 幻灯片列表，每项格式：
            {"type": "content", "title": str, "content": [str, ...]}  — 标准内容页
            {"type": "section", "title": str}                         — 章节分隔页
            {"type": "two_column", "title": str, "left": [...], "right": [...]} — 双栏页
        output_path: 输出文件路径

    Returns:
        str: 生成的 PPTX 路径
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 封面
    _add_title_slide(prs, title)

    # 内容幻灯片
    slide_count = 0
    for sd in slides_data:
        slide_type = sd.get("type", "content")
        if slide_type == "section":
            _add_section_slide(prs, sd.get("title", ""))
        elif slide_type == "two_column":
            _add_two_column_slide(prs, sd.get("title", ""),
                                  sd.get("left", []), sd.get("right", []))
        else:
            _add_content_slide(prs, sd.get("title", ""), sd.get("content", []))
        slide_count += 1

    prs.save(output_path)
    logger.info(f"专业级 PPT 已生成：{output_path} ({slide_count + 1} 页)")
    return output_path


def generate_from_text(text, output_path):
    """将纯文本解析为专业级 PPTX

    支持的格式：

        演示标题
        ## 章节标题
        - 要点1
        - 要点2

        幻灯片标题
        内容段落
        - 子要点

        左栏内容 | 右栏内容

    Args:
        text: 输入文本
        output_path: 输出路径

    Returns:
        str or None: 成功返回路径，失败返回 None
    """
    parsed = _parse_text_to_slides(text)
    if not parsed:
        return None

    title = parsed["title"]
    slides = parsed["slides"]

    return generate_presentation(title, slides, output_path)


# ============================================================
# Pandoc 专业管道（可选）
# ============================================================

def check_pandoc():
    """检查 Pandoc 是否可用"""
    try:
        result = subprocess.run(["pandoc", "--version"], capture_output=True, timeout=5)
        if result.returncode == 0:
            version = result.stdout.decode("utf-8").split("\n")[0]
            return version
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None


def _text_to_markdown(title, slides):
    """将结构化幻灯片数据转换为 Pandoc 兼容 Markdown

    格式：
        % 标题
        # 章节
        ## 幻灯片标题
        - 要点
    """
    lines = [f"% {title}", ""]
    for sd in slides:
        slide_type = sd.get("type", "content")
        slide_title = sd.get("title", "")
        if slide_type == "section":
            lines.append(f"# {slide_title}")
        else:
            lines.append(f"## {slide_title}")
            content = sd.get("content", [])
            if slide_type == "two_column":
                left = sd.get("left", [])
                right = sd.get("right", [])
                for i in range(max(len(left), len(right))):
                    l = left[i] if i < len(left) else ""
                    r = right[i] if i < len(right) else ""
                    if l or r:
                        lines.append(f"- {l} | {r}")
            else:
                for item in content:
                    if "：" in item or ":" in item:
                        lines.append(f"- **{item.split('：')[0]}**：{item.split('：', 1)[1] if len(item.split('：')) > 1 else ''}")
                    else:
                        lines.append(f"- {item}")
        lines.append("")
    return "\n".join(lines)


def generate_via_pandoc(text, output_path, template_path=None):
    """使用 Pandoc 生成专业级 PPTX（需提前安装 pandoc）

    Args:
        text: 输入文本
        output_path: 输出路径
        template_path: 参考模板路径（可选）

    Returns:
        str or None: 成功返回路径，失败返回 None
    """
    pandoc_version = check_pandoc()
    if not pandoc_version:
        logger.warning("Pandoc 未安装，回退到 python-pptx 引擎")
        return None

    parsed = _parse_text_to_slides(text)
    if not parsed:
        return None

    md_content = _text_to_markdown(parsed["title"], parsed["slides"])
    md_path = output_path + ".md"
    Path(md_path).write_text(md_content, encoding="utf-8")

    try:
        cmd = ["pandoc", md_path, "-o", output_path, "--slide-level=2"]
        if template_path and Path(template_path).exists():
            cmd.extend(["--reference-doc", template_path])
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        if result.returncode == 0:
            logger.info(f"Pandoc PPT 已生成：{output_path}")
            Path(md_path).unlink(missing_ok=True)
            return output_path
        else:
            logger.error(f"Pandoc 失败: {result.stderr.decode()}")
            Path(md_path).unlink(missing_ok=True)
            return None
    except Exception as e:
        logger.error(f"Pandoc 异常: {e}")
        Path(md_path).unlink(missing_ok=True)
        return None


if __name__ == "__main__":
    # 测试生成
    test_text = """2026年度工作总结

## 项目进展

成功交付3个核心项目
团队规模扩展至15人
客户满意度提升至95%

技术架构升级
- 完成微服务改造
- 部署自动化CI/CD
- 引入AI辅助开发

## 团队建设

- 新增5名高级工程师
- 建立技术分享机制
- 推行代码评审制度

人才培养 | 技术提升
内部培训12场 | 新技术框架
导师制度 | 性能优化
知识库建设 | 质量改进

## 2027年规划

- 深化AI应用场景
- 拓展海外市场
- 构建数据中台"""

    output = "/tmp/test_professional.pptx"
    result = generate_from_text(test_text, output)
    print(f"✅ 专业级 PPT 测试生成: {result}" if result else "❌ 生成失败")

    pandoc_v = check_pandoc()
    print(f"Pandoc: {pandoc_v if pandoc_v else '未安装（推荐安装以获得最佳效果）'}")
    print(f"安装 Pandoc: brew install pandoc")
